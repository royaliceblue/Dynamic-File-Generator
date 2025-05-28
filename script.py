#!/usr/bin/env python3
import argparse
import os
import re
import io
import zipfile
from zipfile import ZipInfo, ZIP_STORED, ZIP_DEFLATED

# Office libraries
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

# PST via Aspose.Email-for-Python-via-NET
from aspose.email.storage.pst import PersonalStorage, FileFormatVersion

# Regex to parse sizes like '150KB', '2.5MB', or plain number (MB default)
_UNIT_RE = re.compile(r'^(?P<value>\d+(\.\d+)?)(?P<unit>KB|MB)?$', re.IGNORECASE)

def parse_args():
    p = argparse.ArgumentParser(
        description="Generate a valid dummy file of arbitrary size (KB or MB)"
    )
    p.add_argument(
        "-f", "--format",
        choices=["docx","xlsx","pptx","pdf","pst","zip"],
        default="docx",
        help="Format to produce"
    )
    p.add_argument(
        "size",
        help="Target size (e.g. 150KB, 2.5MB, or 10 for 10MB)"
    )
    p.add_argument(
        "output",
        nargs="?",
        help="Output filename (default: <format>_<size>.<ext>)"
    )
    # Allow mixing flags and positionals
    if hasattr(p, "parse_intermixed_args"):
        return p.parse_intermixed_args()
    return p.parse_args()

def parse_size(size_str: str) -> int:
    m = _UNIT_RE.match(size_str.strip())
    if not m:
        raise argparse.ArgumentTypeError(
            f"Invalid size '{size_str}'. Use e.g. 150KB, 2.5MB, or 10"
        )
    val = float(m.group("value"))
    unit = (m.group("unit") or "MB").upper()
    return int(val * (1024 if unit == "KB" else 1024*1024))

def _embed_pad_in_zip(zip_path: str, media_folder: str, target_bytes: int):
    """
    Rewrite the OOXML ZIP at zip_path, patch [Content_Types].xml once,
    then add a pad.bin entry under media_folder so final size == target_bytes.
    """
    # 1) Read original entries
    with zipfile.ZipFile(zip_path, "r") as zf:
        infos = zf.infolist()
        entries = [(info, zf.read(info.filename)) for info in infos]

    # 2) Build a new ZIP in memory
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as newz:
        for info, data in entries:
            # Patch [Content_Types].xml once
            if info.filename == "[Content_Types].xml":
                xml = data.decode("utf-8")
                if 'Extension="bin"' not in xml:
                    xml = xml.replace(
                        "</Types>",
                        '  <Default Extension="bin" ContentType="application/octet-stream"/>\n</Types>'
                    )
                    data = xml.encode("utf-8")
            # Preserve original compression
            new_info = ZipInfo(info.filename)
            new_info.compress_type = info.compress_type
            newz.writestr(new_info, data, compress_type=info.compress_type)

        # 3) Embed pad.bin to reach target size
        current = buf.getbuffer().nbytes
        pad_bytes = target_bytes - current
        if pad_bytes < 0:
            raise ValueError(f"{zip_path!r} is already larger than {target_bytes} bytes")

        pad_name = f"{media_folder}/pad.bin" if media_folder else "pad.bin"
        pad_info = ZipInfo(pad_name)
        pad_info.compress_type = ZIP_STORED
        with newz.open(pad_info, "w") as f:
            # write in 1 MB chunks
            chunk = 1024 * 1024
            full, rem = divmod(pad_bytes, chunk)
            for _ in range(full):
                f.write(b"\0" * chunk)
            if rem:
                f.write(b"\0" * rem)

    # 4) Overwrite the original file with our in-memory ZIP
    with open(zip_path, "wb") as f:
        f.write(buf.getvalue())

def pad_file_trailer(path: str, target_bytes: int):
    """
    For non-ZIP formats (PDF, PST), just append zeros after EOF.
    """
    current = os.path.getsize(path)
    pad_bytes = target_bytes - current
    if pad_bytes < 0:
        raise ValueError(f"{path!r} is already {current} bytes, exceeds target {target_bytes}")
    with open(path, "ab") as f:
        f.write(b"\0" * pad_bytes)

def generate_docx(target_bytes: int, output: str):
    doc = Document()
    doc.add_paragraph(" ")
    doc.save(output)
    _embed_pad_in_zip(output, "word/media", target_bytes)

def generate_xlsx(target_bytes: int, output: str):
    wb = Workbook()
    wb.active.title = "Sheet1"
    wb.save(output)
    _embed_pad_in_zip(output, "xl/media", target_bytes)

def generate_pptx(target_bytes: int, output: str):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(output)
    _embed_pad_in_zip(output, "ppt/media", target_bytes)

def generate_pdf(target_bytes: int, output: str):
    c = canvas.Canvas(output)
    c.drawString(100, 750, "")
    c.showPage()
    c.save()
    pad_file_trailer(output, target_bytes)

def generate_pst(target_bytes: int, output: str):
    with PersonalStorage.create(output, FileFormatVersion.UNICODE) as pst:
        pst.root_folder.add_sub_folder("Inbox")
    pad_file_trailer(output, target_bytes)

def generate_zip(target_bytes: int, output: str):
    # Create minimal ZIP stub
    with zipfile.ZipFile(output, "w", compression=ZIP_STORED) as z:
        z.writestr("dummy.txt", "placeholder")
    # Then embed padding inside the archive root
    _embed_pad_in_zip(output, "", target_bytes)

def main():
    args = parse_args()
    target_bytes = parse_size(args.size)
    fmt = args.format
    out = args.output or f"{fmt}_{args.size.lower()}.{fmt}"

    dispatch = {
        "docx": generate_docx,
        "xlsx": generate_xlsx,
        "pptx": generate_pptx,
        "pdf":  generate_pdf,
        "pst":  generate_pst,
        "zip":  generate_zip,
    }
    dispatch[fmt](target_bytes, out)

    actual = os.path.getsize(out)
    print(f"✅ Generated '{out}' ({actual} bytes, ≈{actual/1024/1024:.2f} MB)")

if __name__ == "__main__":
    main()
